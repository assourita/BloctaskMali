#!/usr/bin/env python
"""
Script complet pour générer la présentation PowerPoint de soutenance BlockTask Mali
Basé sur le mémoire de fin d'études Master - 42 diapositives
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_blocktask_presentation():
    """Crée la présentation PowerPoint complète de soutenance BlockTask"""
    
    # Création de la présentation
    prs = Presentation()
    
    # Styles
    title_font_size = Pt(36)
    subtitle_font_size = Pt(24)
    content_font_size = Pt(18)
    bullet_font_size = Pt(16)
    
    # Couleurs BlockTask
    primary_color = RGBColor(0, 123, 255)  # Bleu
    secondary_color = RGBColor(40, 167, 69)  # Vert
    accent_color = RGBColor(255, 193, 7)  # Jaune
    
    # Diapo 1: Titre du projet / Qui suis-je
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "BlockTask Mali"
    subtitle.text = "Plateforme hybride de délégation sécurisée de tâches physiques\n\nMémoire de fin d'études Master\nILSI - Informatique et Logiciels pour les Systèmes d'Information\n\nPrésenté par: [Votre Nom]\nEncadrant: Dr Oumar MAIGA\nITMA - [Année]"
    
    title.text_frame.paragraphs[0].font.size = title_font_size
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    subtitle.text_frame.paragraphs[0].font.size = subtitle_font_size
    
    # Diapo 2: Chapitre 1 - Contextes
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 1: CONTEXTES"
    content.text = """Trois niveaux de contexte analysés:

• Contexte général: Économie numérique et plateformes de services
• Contexte africain et ouest-africain: Mobile Money et inclusion financière  
• Contexte national: Le Mali - Marché cible et spécificités locales
• Contexte technologique: Architecture hybride BlockTask"""
    
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Diapo 3: Contexte général
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexte Général"
    content.text = """Transformation numérique et économie à la tâche (gig economy):

• Plateformes numériques: mise en relation, coordination, gestion confiance
• Trois piliers: mise en relation rapide, coordination des opérations, confiance
• Critiques des plateformes centralisées: 
  - Opacité des commissions
  - Dépendance vis-à-vis d'intermédiaires
  - Litiges difficiles à trancher
  - Exclusion des acteurs peu bancarisés"""
    
    # Diapo 4: Contexte africain
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexte Africain et Ouest-Africain"
    content.text = """Croissance du Mobile Money comme vecteur d'inclusion financière:

• Espace UEMOA: Franc CFA (XOF) comme monnaie commune
• Opérateurs dominants: Orange Money et Moov Money
• Obstacles récurrents:
  - Manque de confiance entre acteurs
  - Absence de mécanismes d'escrow accessibles
  - Difficulté de preuve d'exécution
  - Fragmentation des moyens de paiement
  - Faible maturité des cadres réglementaires"""
    
    # Diapo 5: Contexte malien
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexte National: Le Mali"
    content.text = """Marché cible de la première phase BlockTask:

• Inclusion numérique: Taux de pénétration mobile significatif
• Mobile Money: Adoption Orange Money et Moov Money
• Paiement en FCFA via Mobile Money = habitudes réelles utilisateurs
• Identité: NINA comme référentiel pour KYC adapté au marché malien
• Géographie: Zones urbaines et périurbaines (Bamako, Ségou, Sikasso)
• Besoins: Livraison, assistance, courses dans les villes"""
    
    # Diapo 6: Contexte technologique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexte Technologique BlockTask"
    content.text = """Architecture hybride innovante:

• Couche Métier: Angular + Django REST (gestion missions, profils, KYC, litiges)
• Paiement Réel: Orange Money / Moov Money (FCFA) - Paiements quotidiens
• Couche Confiance: Blockchain Ethereum (testnet Sepolia) - Ancrage optionnel
• Identité: NINA + validation KYC admin - Conformité et sécurité

Objectif: Concilier usage local Mobile Money + apports blockchain traçabilité"""
    
    # Diapo 7: Chapitre 2 - Problématique (Constats)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 2: PROBLÉMATIQUE - Constats"
    content.text = """Difficultés récurrentes au Mali:

1. Asymétrie d'information: Client ne peut vérifier fiabilité prestataire
2. Risque de non-exécution: Absence garantie tâche réalisée correctement
3. Risque de non-paiement: Prestataire craint non-rémunération
4. Gestion litiges: Pas mécanisme structuré médiation/arbitrage
5. Inadéquation solutions existantes: 
   - Plateformes internationales: Carte bancaire/PayPal inadaptés
   - Solutions blockchain: Excluent utilisateurs non-technophiles"""
    
    # Diapo 8: Question de recherche
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Question de Recherche"
    content.text = """Problématique centrale:

Comment concevoir et mettre en œuvre une plateforme numérique de délégation de tâches, adaptée au contexte malien (FCFA, Mobile Money, NINA), intégrant de manière hybride la blockchain pour renforcer la confiance transactionnelle, tout en restant utilisable par des acteurs peu technophiles?

Sous-questions:
• Q1: Concepts et travaux pertinents pour contexte malien?
• Q2: Architecture efficace couplant backend, Mobile Money, blockchain?
• Q3: Mécanismes sécurité, identification, preuve d'exécution nécessaires?
• Q4: Prototype répond-il aux objectifs et contraintes marché malien?"""
    
    # Diapo 9: Chapitre 3 - Objectifs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 3: OBJECTIFS"
    content.text = """Objectif général:
Concevoir, développer et évaluer une plateforme web hybride de délégation sécurisée de tâches physiques, adaptée au marché malien, combinant paiements Mobile Money en FCFA et ancrage blockchain via contrats intelligents.

Objectifs spécifiques:
• OS1: Analyser domaine et état de l'art
• OS2: Spécifier besoins fonctionnels et non fonctionnels  
• OS3: Concevoir architecture modulaire et évolutive
• OS4: Implémenter espaces utilisateurs (client, prestataire, entreprise, admin)
• OS5: Intégrer paiements Mobile Money
• OS6: Intégrer couche blockchain
• OS7: Mettre œuvre sécurité et conformité
• OS8: Réaliser suivi et preuves d'exécution
• OS9: Tester et valider système
• OS10: Évaluer atteinte objectifs"""
    
    # Diapo 10: Hypothèses de travail
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Hypothèses de Travail"
    content.text = """H1: Modèle hybride (Mobile Money + blockchain optionnelle) plus adapté au Mali qu'une solution exclusivement crypto ou exclusivement centralisée.

H2: Smart contracts d'escrow améliorent perception confiance et traçabilité, même si paiement réel reste en FCFA.

H3: Architecture monolithique modulaire (Django) permet de livrer prototype cohérent dans délai mémoire Master, tout en restant évolutive."""
    
    # Diapo 11: Méthodologie - Approche globale
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 4: MÉTHODOLOGIE - Approche Globale"
    content.text = """Démarche inspirée cycle de vie itératif:

Analyse du domaine → État de l'art → Conception UML/C4 → Implémentation → Tests → Discussion et conclusion

Approche par prototypage incrémental:
• Chaque itération enrichit le système
• Version démontrable à chaque étape
• Authentification → Missions → Paiements → Blockchain

Méthodes et techniques:
• Analyse et conception: UML, C4, identification acteurs/processus
• Développement: Architecture 3-tiers, services métier, API REST
• Qualité: Tests fonctionnels, automatisés, matrice traçabilité
• Sécurité: Analyse menaces, processus KYC, bonnes pratiques"""
    
    # Diapo 12: Outils et environnement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Outils et Environnement"
    content.text = """Stack technique principal:

Backend: Python/Django REST Framework, PostgreSQL, Redis
Frontend: TypeScript/Angular 17, Material Design
Mobile: React Native/Expo
Blockchain: Solidity, Hardhat, testnet Sepolia
Tests: pytest, Hardhat + Chai, scénarios manuels
Déploiement: Docker, Render (production), Git (versionnement)

Outils additionnels:
• Django Channels (WebSocket temps réel)
• Celery (tâches asynchrones) 
• GitHub Actions (intégration continue)
• OpenZeppelin (sécurité smart contracts)"""
    
    # Diapo 13: Chapitre 5 - État de l'art Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 5: ÉTAT DE L'ART"
    content.text = """Revue complète des fondements et solutions existantes:

• Concepts fondamentaux: plateformes services, escrow, réputation, blockchain
• Analyse plateformes centralisées référence: Upwork, Fiverr, TaskRabbit
• Systèmes réputation numériques: mécanismes, vulnérabilités, pistes blockchain
• Blockchain et smart contracts: fondements, apports pour délégation tâches
• Limites approches purement blockchain: complexité, adoption, coûts
• Plateformes décentralisées: Ethlance, Braintrust, travaux académiques
• Positionnement BlockTask: avantages hybride vs solutions existantes"""
    
    # Diapo 14: Concepts fondamentaux (Partie 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Concepts Fondamentaux (1/2)"
    content.text = """Économie à la tâche et plateformes de mise en relation:
• Transformation numérique modes organisation travail
• Plateformes: marketplaces, mobilité, micro-tâches
• Trois piliers: mise en relation, coordination, confiance

Coûts de transaction et rôle de l'escrow:
• Théorie coûts transaction (Coase, Williamson)
• Escrow comme mécanisme réduction asymétrie information
• Séquestre fonds jusqu'à validation preuve exécution

Asymétrie d'information et systèmes de réputation:
• Problème principal: information incomplète entre acteurs
• Réputation comme signal qualité et fiabilité
• Vulnérabilités: sybil attacks, review bombing

Preuve d'exécution et tâches physiques:
• Difficulté attestation réalisation services physiques
• Photos, GPS, signatures comme mécanismes preuve"""
    
    # Diapo 15: Concepts fondamentaux (Partie 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Concepts Fondamentaux (2/2)"
    content.text = """Blockchain, registre distribué et contrats intelligents:
• Ethereum: machine virtuelle et smart contracts
• Immuabilité, transparence, exécution automatisée
• Solidity: langage développement contrats

Architectures hybrides (off-chain / on-chain):
• Séparation traitement métier et ancrage confiance
• Off-chain: rapidité, coût faible, expérience utilisateur
• On-chain: traçabilité, sécurité, engagement contractuel

Mobile Money, inclusion financière et contexte ouest-africain:
• GSMA: 1+ milliard comptes Mobile Money Afrique
• XOF comme monnaie commune UEMOA  
• Orange Money, Moov Money: opérateurs dominants

KYC et identité numérique (NINA):
• Numéro Identification Nationale Mali comme référentiel
• Processus Know Your Customer adapté contexte local"""
    
    # Diapo 16: Plateformes centralisées de référence
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Plateformes Centralisées de Référence"
    content.text = """Upwork:
• Leader mondial freelancing informatique
• Système réputation sophistiqué (Job Success Score)
• Escrow intégré, protection travailleurs
• Limites: focus digital, paiement carte bancaire

Fiverr:
• Marketplace services créatifs et digitaux
• Modèle "gigs" (services prédéfinis)
• Système niveaux (New, Level 1, Level 2, Top Rated)
• Limites: principalement services digitaux

TaskRabbit:
• Focus services physiques locaux (bricolage, livraison)
• Modèle basé localisation géographique
• "Taskers" avec background checks
• Limites: principalement marché nord-américain"""
    
    # Diapo 17: Synthèse comparative plateformes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Synthèse Comparative Plateformes"
    content.text = """Comparaison Upwork/Fiverr/TaskRabbit:

Forces:
• Maturité opérationnelle et volume transactions
• Systèmes réputation éprouvés
• Expérience utilisateur optimisée
• Protection escrow intégrée

Faiblesses:
• Inadaptation moyens paiement locaux (FCFA, Mobile Money)
• Modèles économiques centrés marchés développés
• Exclusion utilisateurs peu bancarisés
• Cadres juridiques étrangers contexte malien

Positionnement BlockTask:
• Spécifiquement conçu pour marché malien/UEMOA
• Intégration Mobile Money + KYC NINA
• Architecture hybride adaptable contexte local"""
    
    # Diapo 18: Chapitre 6 - Conception Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CHAPITRE 6: CONCEPTION DE LA SOLUTION"
    content.text = """Introduction:

Analyse structurée du domaine et modélisation complète:
• Périmètre fonctionnel et acteurs du système
• Processus métier principal: cycle de vie mission
• Règles de gestion et exigences
• Modélisation architecturale C4 (4 niveaux)
• Modélisation UML complète (cas d'utilisation, classes, séquences)
• Scénarios d'utilisation détaillés
• Matrice de traçabilité conception → implémentation"""
    
    # Diapo 19: Périmètre fonctionnel et acteurs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Périmètre Fonctionnel et Acteurs"
    content.text = """Périmètre fonctionnel BlockTask:
• Création et gestion missions
• Processus de candidature et sélection
• Gestion paiements et escrow
• Suivi exécution et preuves
• Système réputation et litiges
• Administration plateforme

Acteurs du système:
• Client: Créateur missions, payeur final
• Prestataire: Exécutant missions, fournisseur services
• Entreprise (B2B): Structure employant prestataires
• Administrateur: Validation KYC, arbitrage litiges, configuration"""
    
    # Diapo 20: Processus métier principal
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Processus Métier: Cycle de Vie Mission"
    content.text = """États principaux du cycle de vie:

1. DRAFT: Brouillon mission (client)
2. PUBLISHED: Publiée (visible prestataires)
3. FUNDED: Financée (escrow activé)
4. ACCEPTED: Acceptée (prestataire sélectionné)
5. DEPOSIT_PAID: Caution déposée (prestataire)
6. IN_PROGRESS: En cours (exécution)
7. COMPLETED: Terminée (preuves soumises)
8. REVIEW: En révision (validation client)
9. PAID: Payée (fonds libérés)
10. DISPUTED: Litige (arbitrage requis)
11. CANCELLED: Annulée (remboursement)
12. EXPIRED: Expirée (délai dépassé)"""
    
    # Diapo 21: Règles de gestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Règles de Gestion (Extrait)"
    content.text = """RG01: Une mission doit être financée pour devenir visible aux prestataires

RG02: Après acceptation candidature, prestataire a 4h pour déposer caution

RG03: La caution est calculée dynamiquement: 10-30% budget mission

RG04: Les fonds sont bloqués en escrow jusqu'à validation finale

RG05: Commission plateforme: 5% du montant mission (client)

RG06: L'accès plateforme est conditionné au statut KYC vérifié

RG07: Le système réputation combine: 40% succès, 30% note, 20% litiges, 10% volume

RG08: Un litige ne peut être ouvert qu'après soumission preuves

RG09: La décision arbitre est finale et impacte réputation parties"""
    
    # Diapo 22: Exigences fonctionnelles
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Exigences Fonctionnelles Principales"
    content.text = """EF01: Inscription/connexion avec gestion multi-rôles (client/prestataire)

EF02: Processus KYC avec NINA et validation administrative

EF03: Création mission avec formulaire guidé et devis estimé

EF04: Paiement escrow via Mobile Money (Orange/Moov)

EF05: Consultation missions disponibles avec filtres et recherche

EF06: Candidature prestataire avec message et portfolio

EF07: Acceptation candidature et notification automatique

EF08: Dépôt caution prestataire (dynamique selon mission)

EF09: Démarrage mission, suivi GPS, soumission preuves

EF10: Validation finale et libération fonds automatique

EF11: Gestion litiges avec arbitrage administrateur

EF12: Ancrage optionnel engagements sur blockchain

EF13: Gestion entreprises B2B et affectation employés

EF14: Tableau bord administration avec supervision complète"""
    
    # Diapo 23: Exigences non fonctionnelles
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Exigences Non Fonctionnelles"
    content.text = """ENF01: Authentification JWT sécurisée avec tokens expirants

ENF02: KYC obligatoire bloquant accès non-vérifiés

ENF03: Pagination API (20 résultats par page)

ENF04: Architecture découplée frontend/backend (CORS)

ENF05: Traçabilité complète actions (historique, logs)

ENF06: Localisation Mali (XOF, +223, villes principales)

ENF07: Modularité (12 apps Django indépendantes)

ENF08: Utilisabilité interface Material Design responsive

ENF09: Documentation API OpenAPI/Swagger accessible

ENF10: Maintenabilité code structuré et commenté

ENF11: Performance temps réponse < 2s

ENF12: Sécurité données (HTTPS, hash mots de passe)"""
    
    # Diapo 24: Modélisation C4 - Définition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Modélisation Architecturale - Modèle C4"
    content.text = """Modèle C4: Architecture logicielle à 4 niveaux

Niveau 1 - Contexte: Système dans environnement global
• Acteurs externes et interactions avec BlockTask

Niveau 2 - Conteneurs: Applications et technologies
• Frontend web, backend, base données, blockchain

Niveau 3 - Composants: Architecture interne backend
• Modules Django, services, API, modèles

Niveau 4 - Code: Détails implémentation (non présenté)

Avantages:
• Vue progressive adaptée à chaque audience
• Communication efficace architecture technique
• Documentation vivante du système"""
    
    # Diapo 25: Diagramme de contexte C4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Contexte C4"
    content.text = """BlockTask au centre de l'écosystème:

Acteurs externes:
• Client: Crée missions, paie, valide
• Prestataire: Postule, exécute, reçoit paiement  
• Entreprise: Gère employés, affecte missions
• Administrateur: Valide KYC, arbitre, configure

Systèmes externes:
• Orange Money API: Paiements Mobile Money
• Moov Money API: Alternative paiements
• Ethereum Blockchain: Ancrage smart contracts
• Email Service: Notifications transactionnelles

Flux principaux:
• Clients ↔ BlockTask: Gestion missions
• Prestataires ↔ BlockTask: Candidatures, exécution
• BlockTask ↔ Mobile Money: Transactions financières
• BlockTask ↔ Blockchain: Ancrage confiance"""
    
    # Diapo 26: Diagramme de conteneurs C4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Conteneurs C4"
    content.text = """Architecture 3-tiers BlockTask:

Frontend (Angular 17):
• Application web monopage (SPA)
• Interface responsive Material Design
• Gestion multi-rôles (client/prestataire/entreprise/admin)
• Communication API REST HTTP/HTTPS

Backend (Django REST):
• 12 applications modulaires
• API REST documentée Swagger
• Authentification JWT, autorisation par rôles
• Logique métier, règles de gestion

Base données (PostgreSQL):
• Données structurées métier
• Migrations Django versionnées
• Backup et récupération

Cache (Redis):
• Sessions utilisateurs, temps réel
• WebSocket pour suivi GPS

Blockchain (Ethereum):
• Smart contracts Solidity
• Réseau testnet Sepolia
• Ancrage optionnel engagements"""
    
    # Diapo 27: Diagramme de composants C4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Composants - Backend Django"
    content.text = """Architecture modulaire backend Django:

Core Apps:
• users: Authentification, profils, KYC
• missions: Gestion missions, workflow
• enterprises: B2B, employés, affectations
• payments: Escrow, Mobile Money intégration
• notifications: Email, WebSocket temps réel
• disputes: Litiges, arbitrage
• reputation: Scores, calculs automatiques

Support Apps:
• blockchain: Smart contracts, ancrage on-chain
• common: Utils, shared models
• analytics: Statistiques, tableaux bord
• admin: Administration avancée

API Layer:
• Viewsets DRF avec permissions granulaires
• Services métier (business logic)
• Serializers validation et transformation"""
    
    # Diapo 28: Modélisation UML Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Modélisation UML Complète"
    content.text = """Diagrams UML pour spécification détaillée:

Diagrammes statiques:
• Cas d'utilisation: Interactions acteurs/système
• Classes: Modèle de données métier

Diagrammes dynamiques:
• Séquences: Flux temporels interactions
• Activités: Processus métier workflow

Diagrammes implémentation:
• Déploiement: Architecture physique production

Objectifs:
• Spécification fonctionnelle précise
• Communication équipe développement
• Base documentation technique
• Référence pour tests et validation"""
    
    # Diapo 29: Cas d'utilisation - Client
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Cas d'Utilisation - Client"
    content.text = """Cas d'utilisation principaux client:

Gestion missions:
• UC-C01: Créer mission (formulaire guidé)
• UC-C02: Consulter missions créées
• UC-C03: Modifier mission (brouillon)
• UC-C04: Annuler mission (remboursement)

Processus recrutement:
• UC-C05: Consulter candidatures
• UC-C06: Accepter candidature prestataire  
• UC-C07: Refuser candidature

Suivi et validation:
• UC-C08: Suivre progression mission (GPS)
• UC-C09: Consulter preuves exécution
• UC-C10: Valider travail final
• UC-C11: Ouvrir litige (si nécessaire)

Gestion compte:
• UC-C12: Gérer profil et préférences
• UC-C13: Consulter historique missions"""
    
    # Diapo 30: Cas d'utilisation - Prestataire
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Cas d'Utilisation - Prestataire"
    content.text = """Cas d'utilisation principaux prestataire:

Recherche missions:
• UC-P01: Consulter missions disponibles
• UC-P02: Rechercher missions (filtres)
• UC-P03: Voir détails mission

Candidature:
• UC-P04: Postuler à mission
• UC-P05: Gérer candidatures envoyées
• UC-P06: Consulter statut candidature

Exécution:
• UC-P07: Accepter mission assignée
• UC-P08: Déposer caution (si requis)
• UC-P09: Démarrer mission
• UC-P10: Soumettre preuves exécution
• UC-P11: Activer suivi GPS

Gestion compte:
• UC-P12: Gérer profil et portfolio
• UC-P13: Consulter réputation et scores
• UC-P14: Gérer paramètres disponibilité"""
    
    # Diapo 31: Cas d'utilisation - Admin et Entreprise
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Cas d'Utilisation - Admin et Entreprise"
    content.text = """Administrateur:

Supervision (UC-A01):
• Valider dossiers KYC (NINA)
• Arbitrer litiges missions
• Configurer paramètres plateforme
• Consulter statistiques et analytics
• Gérer utilisateurs et permissions

Entreprise B2B:

Gestion employés (UC-E01):
• Créer et gérer comptes employés
• Affecter employés aux missions
• Consulter disponibilités équipe
• Valider temps travail employés

Supervision (UC-E02):
• Consulter missions entreprise
• Suivre performance équipe
• Gérer facturation et paiements
• Exporter rapports activité"""
    
    # Diapo 32: Diagramme de classes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Classes - Cœur Métier"
    content.text = """Modèle de données principal:

User (AbstractUser):
• email, username, first_name, last_name
• is_active, is_suspended, user_type
• Méthodes: get_full_name(), has_role()

Profile (OneToOne User):
• phone, address, city, country
• avatar, bio, preferences
• KYC: nina, kyc_status, kyc_documents

Mission:
• title, description, budget, currency
• status, deposit_required, deposit_paid
• client (FK User), provider (FK User)
• executing_employee (FK Employee)

Employee (Profile):
• enterprise (FK Enterprise), position
• is_active, terminated_at
• user (OneToOne User)

Enterprise:
• company_name, company_email, phone
• address, business_license
• user (FK User)

Models support:
• MissionApplication, EscrowTransaction
• Dispute, Reputation, Notification"""
    
    # Diapo 33: Diagramme de séquence - Création mission
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Séquence - Création Mission et Paiement"
    content.text = """Flux création mission et paiement escrow:

1. Client → Frontend: Remplir formulaire mission
2. Frontend → API: POST /api/missions/
3. API → Validation: Données requises + KYC client
4. API → Database: Créer mission (status: DRAFT)
5. API → Client: Mission créée (ID généré)

6. Client → Frontend: "Payer maintenant"
7. Frontend → Mobile Money: Initier paiement
8. Mobile Money → Client: OTP validation (1234)
9. Mobile Money → API: Confirmation paiement
10. API → Escrow Service: Bloquer fonds
11. API → Database: Mission status = FUNDED
12. API → Notifications: Email confirmation
13. API → Blockchain (optionnel): Record mission

Résultat: Mission visible prestataires, fonds sécurisés"""
    
    # Diapo 34: Diagramme de séquence - Exécution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Séquence - Exécution et Validation"
    content.text = """Flux exécution mission et validation:

1. Prestataire → API: POST /api/missions/{id}/start/
2. API → Validation: Permissions + deposit_paid
3. API → Database: Mission status = IN_PROGRESS
4. API → WebSocket: Démarrer suivi GPS

5. Prestataire → API: Upload preuves (photos, notes)
6. API → Storage: Sauvegarder fichiers
7. API → Database: Preuves enregistrées

8. Prestataire → API: POST /api/missions/{id}/complete/
9. API → Database: Mission status = COMPLETED
10. API → Client: Notification validation requise

11. Client → API: POST /api/missions/{id}/validate/
12. API → Escrow Service: Libérer fonds (95% prestataire, 5% commission)
13. API → Reputation: Mettre à jour scores
14. API → Blockchain (optionnel): Record completion

Résultat: Mission terminée, paiements effectués"""
    
    # Diapo 35: Diagramme d'activité
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme d'Activité - Processus Global"
    content.text = """Workflow mission complet:

[Client] → Créer mission → [Formulaire]
                ↓
[Payer escrow] → [Mobile Money] → [Fonds bloqués]
                ↓
[Publication] → [Visible prestataires]

[Prestataire] → [Postuler] → [Candidature]
                ↓
[Client] → [Examiner] → [Accepter/Refuser]
                ↓
[Prestataire] → [Déposer caution] → [4h délai]
                ↓
[Démarrer] → [Suivi GPS] → [Exécuter]
                ↓
[Soumettre preuves] → [Photos + notes]
                ↓
[Client] → [Valider] → [Libérer fonds]
                ↓
[Mission terminée] → [Mise à jour réputation]

[Gestion litiges] → [Arbitrage admin] → [Décision finale]"""
    
    # Diapo 36: Diagramme de déploiement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Diagramme de Déploiement"
    content.text = """Architecture production prototype:

Infrastructure Cloud (Render):
• Frontend: Angular build statique (CDN)
• Backend: Django REST API (conteneur Docker)
• Base données: PostgreSQL géré
• Cache: Redis géré

Réseau externe:
• Mobile Money: APIs Orange/Moov (HTTPS)
• Blockchain: Ethereum Sepolia (RPC)
• Email: Service transactionnel (SMTP)

Développement local:
• Docker Compose: environnement complet
• Hot reload: frontend/backend
• Base SQLite: développement rapide
• Hardhat local: blockchain test

Sécurité:
• HTTPS/TLS 1.3
• Firewall cloud
• Secrets management (variables environnement)
• Backup automatique base données"""
    
    # Diapo 37: Scénarios détaillés (Partie 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Scénarios d'Utilisation Détaillés (1/3)"
    content.text = """SC1: Inscription, KYC et connexion

Acteur: Nouvel utilisateur (prestataire)
Préconditions: Accès internet, email valide, NINA

Étapes:
1. Accéder page inscription BlockTask
2. Choisir profil "Prestataire" 
3. Remplir informations personnelles (+223)
4. Soumettre NINA + documents (carte identité)
5. Recevoir email confirmation inscription
6. Statut: "En attente validation KYC"
7. Admin valide dossier dans interface admin
8. Recevoir email "Compte vérifié"
9. Se connecter avec email/mot de passe
10. Accéder tableau bord prestataire

Résultat: Compte actif, missions visibles"""
    
    # Diapo 38: Scénarios détaillés (Partie 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Scénarios d'Utilisation Détaillés (2/3)"
    content.text = """SC2: Création mission et paiement escrow

Acteur: Client (entreprise)
Préconditions: Compte vérifié, solde Mobile Money

Étapes:
1. Connecter espace client
2. Cliquer "Créer nouvelle mission"
3. Remplir formulaire en 4 étapes:
   - Description et budget
   - Localisation et date
   - Compétences requises
   - Confirmation
4. Choisir "Payer avec Orange Money"
5. Recevoir OTP (1234 en sandbox)
6. Valider paiement
7. Confirmation: "Mission publiée et financée"
8. Recevoir candidatures prestataires

Résultat: Mission visible, fonds sécurisés"""
    
    # Diapo 39: Scénarios détaillés (Partie 3)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Scénarios d'Utilisation Détaillés (3/3)"
    content.text = """SC3: Exécution, preuves et validation

Acteurs: Prestataire et Client
Préconditions: Mission acceptée, caution déposée

Étapes prestataire:
1. Démarrer mission (bouton "Commencer")
2. Activer suivi GPS temps réel
3. Exécuter tâche selon cahier charges
4. Prendre photos progression
5. Soumettre preuves finale
6. Marquer mission comme terminée

Étapes client:
1. Recevoir notification "Mission terminée"
2. Consulter preuves soumises
3. Valider si satisfait OU ouvrir litige
4. Confirmer validation finale

Résultat: Paiement automatique, mise à jour réputation"""
    
    # Diapo 40: Tests et résultats obtenus
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Tests et Résultats Obtenus"
    content.text = """Stratégie test multi-niveaux:

Tests automatisés (26 total):
• Backend Django: 11 tests pytest
  - Flux MVP, KYC, rôles, litiges, réputation
• Smart contracts: 15 tests Hardhat
  - Escrow, réputation, litigation (5 tests chacun)

Tests fonctionnels manuels (7 scénarios):
• SC1: Inscription, KYC, connexion ✅
• SC2: Création mission, paiement ✅  
• SC3: Candidature, acceptation ✅
• SC4: Exécution, preuves, GPS ✅
• SC5: Validation, paiement ✅
• SC6: Litiges, arbitrage ✅
• SC7: Administration ✅

Résultats globaux:
• Taux couverture: 80-100% exigences
• Prototype fonctionnel complet
• Objectifs spécifiques: 8/10 atteints, 2/10 partiels"""
    
    # Diapo 41: Perspectives et difficultés
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Perspectives et Difficultés"
    content.text = """Difficultés rencontrées et solutions:

• Inadéquation paiement crypto/habitudes Mali
  → Modèle hybride FCFA + ancrage optionnel

• Adresses contrats changeantes (Remix/Hardhat)  
  → Scripts apply-config, remix:addresses

• Complexité sync Django ↔ blockchain
  → API record-mission, record-proof, sync-events

• Configuration RPC/clés blockchain
  → Documentation README_MALI, guides

Perspectives d'évolution:

Court terme (3-6 mois):
• Intégration API Mobile Money production
• Finalisation APK mobile Expo
• Audit sécurité smart contracts

Moyen terme (6-18 mois):
• Extension UEMOA (Sénégal, Côte d'Ivoire)
• Oracles Chainlink preuves on-chain
• Étude utilisateur terrain Bamako

Long terme (18+ mois):
• Passage mainnet après audit
• DAO gouvernance plateforme
• Partenariats télécoms/microfinance"""
    
    # Diapo 42: Conclusion et remerciements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion et Remerciements"
    content.text = """Conclusion:

BlockTask démontre qu'il est possible de concilier:
✅ Innovation blockchain et réalités locales africaines
✅ Inclusion financière via Mobile Money  
✅ Confiance institutionnelle via KYC NINA
✅ Traçabilité renforcée via smart contracts (optionnel)

Le Mali, avec sa dynamique Mobile Money et son besoin de structuration de l'économie informelle, constitue un terrain d'expérimentation pertinent.

Ce mémoire pose les fondations techniques et conceptuelles; la suite dépendra des partenariats, du cadre réglementaire et de l'adoption par les utilisateurs.

Remerciements:

• Dr Oumar MAIGA - Encadrement et disponibilité
• ITMA - Formation de qualité et ressources
• Famille et amis - Soutien moral et patience
• Tous ceux qui croient au potentiel du numérique pour l'Afrique

Merci pour votre attention!

Questions?"""
    
    # Sauvegarde de la présentation
    output_path = os.path.join(os.path.dirname(__file__), "BlockTask_Soutenance_Complete.pptx")
    prs.save(output_path)
    
    print(f"✅ Présentation BlockTask générée avec succès!")
    print(f"📍 Fichier sauvegardé: {output_path}")
    print(f"📊 Nombre de diapos: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_blocktask_presentation()
