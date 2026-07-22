from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


LOGO_PATH = r"C:\Users\PC\Documents\blocktask_myself\frontend\src\assets\images\logo-blocktask-mali.png"


def add_title_slide(prs, title, subtitle, logo_path=None):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    if logo_path:
        try:
            slide.shapes.add_picture(
                logo_path,
                Inches(11.0), Inches(0.3),
                width=Inches(1.8)
            )
        except Exception:
            pass
    return slide


def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    return slide


def add_demo_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Bullets en haut à gauche, compactes
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

    # Grand placeholder pour la démo
    left = Inches(1.0)
    top = Inches(3.0)
    width = Inches(11.3)
    height = Inches(3.8)
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    box.line.color.rgb = RGBColor(150, 150, 150)
    box.line.width = Pt(2)

    tf_box = box.text_frame
    tf_box.text = "[ZONE DÉMO : capture d'écran ou live]"
    p = tf_box.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(100, 100, 100)
    tf_box.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def add_content_slide(prs, title, bullets, image_placeholder=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)

    if image_placeholder:
        # Add a placeholder text box for the image/diagram
        left = Inches(8.5)
        top = Inches(2.0)
        width = Inches(4.3)
        height = Inches(4.5)
        box = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.color.rgb = RGBColor(150, 150, 150)
        box.line.width = Pt(2)

        tf_box = box.text_frame
        tf_box.text = image_placeholder
        p = tf_box.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Center vertically
        tf_box.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Garde
    add_title_slide(
        prs,
        "BlockTask",
        "Plateforme hybride de délégation sécurisée de tâches physiques — cas du Mali\n\nSoutenance de mémoire de fin d'études — Master ILSI",
        logo_path=LOGO_PATH
    )

    # 2. Contexte
    add_content_slide(prs, "Contexte et motivation", [
        "Économie à la tâche au Mali : encore largement informelle",
        "Obstacles : manque de confiance, absence d'escrow accessible, paiements internationaux inadaptés",
        "Opportunité : Mobile Money (Orange Money, Moov Money) très répandu",
        "Leviers locaux : FCFA (XOF), NINA, blockchain optionnelle"
    ])

    # 3. Problématique
    add_content_slide(prs, "Problématique", [
        "Comment concevoir et mettre en œuvre une plateforme numérique de délégation de tâches,",
        "adaptée au contexte malien (FCFA, Mobile Money, NINA),",
        "intégrant de manière hybride la blockchain pour renforcer la confiance transactionnelle,",
        "tout en restant utilisable par des acteurs peu technophiles ?"
    ])

    # 4. Objectifs
    add_content_slide(prs, "Objectifs spécifiques", [
        "OS1 — Analyser le domaine et l'état de l'art",
        "OS2/OS3 — Spécifier et concevoir une architecture modulaire (UML / C4)",
        "OS4 — Implémenter les espaces client, prestataire, entreprise, admin",
        "OS5 — Intégrer les paiements Mobile Money en FCFA",
        "OS6 — Intégrer une couche blockchain optionnelle (Solidity / Hardhat)",
        "OS7 — Sécurité, KYC NINA et conformité",
        "OS8 — Suivi GPS et preuves d'exécution",
        "OS9/OS10 — Tester et évaluer le système"
    ])

    # 5. État de l'art
    add_content_slide(prs, "État de l'art et positionnement", [
        "Plateformes centralisées : Upwork, Fiverr, TaskRabbit",
        "Limites : centralisation, commissions élevées, inadaptation au Mali",
        "Plateformes décentralisées : Ethlance, Braintrust — crypto obligatoire",
        "Positionnement BlockTask : hybride — Mobile Money + blockchain optionnelle"
    ])

    # 6. Architecture
    add_content_slide(prs, "Architecture globale", [
        "Web : Angular 17 + Material",
        "Mobile : React Native + Expo SDK 52",
        "Backend : Django 4.2 + DRF + PostgreSQL",
        "Temps réel : Channels + Redis + Celery",
        "Blockchain : Solidity + Hardhat / Sepolia",
        "Déploiement : Render + Railway"
    ], image_placeholder="[INSÉRER : Figure 3.3 — Diagramme de conteneurs C4]")

    # 7. Cycle de vie
    add_content_slide(prs, "Cycle de vie d'une mission", [
        "Création → Paiement escrow (FCFA)",
        "Candidature → Acceptation → Caution",
        "Exécution → Preuves",
        "Validation → Libération fonds (95 % / 5 %)",
        "Litige possible à chaque étape"
    ], image_placeholder="[INSÉRER : Figure 3.1 — Diagramme d'états mission]")

    # 8. Implémentation
    add_content_slide(prs, "Fonctionnalités clés implémentées", [
        "Auth JWT avec rôles doubles",
        "KYC NINA + validation admin",
        "Paiement Mobile Money sandbox",
        "Caution dynamique selon réputation",
        "Chat mission + 2FA TOTP",
        "GPS temps réel + preuves photo",
        "Analyse automatique des preuves"
    ], image_placeholder="[INSÉRER : captures interfaces web / mobile]")

    # 9. Sécurité
    add_content_slide(prs, "Sécurité et confiance", [
        "JWT access 1h / refresh 7j",
        "Guards Angular multi-rôles",
        "KYC bloquant avant accès",
        "OpenZeppelin : ReentrancyGuard, Ownable, Pausable",
        "Escrow FCFA verrouillé",
        "WebSocket GPS authentifié"
    ], image_placeholder="[INSÉRER : schéma flux JWT + KYC]")

    # 10. Déploiement
    add_content_slide(prs, "Déploiement et environnements", [
        "Développement : Docker Compose",
        "Production : Render (backend + frontend)",
        "Redis : WebSockets + Celery",
        "Mobile : APK Android via Gradle",
        "Blockchain : Sepolia + apply-config"
    ], image_placeholder="[INSÉRER : Figure 3.10 — Diagramme de déploiement]")

    # 11. Tests
    add_content_slide(prs, "Tests et validation", [
        "11 tests backend (pytest)",
        "15 tests smart contracts (Hardhat)",
        "7 scénarios manuels",
        "CI GitHub Actions",
        "Matrice exigences ↔ tests",
        "Atteinte objectifs : 80–100 %"
    ], image_placeholder="[INSÉRER : Tableau 5.7 — Synthèse quantitative]")

    # 12. Résultats
    add_content_slide(prs, "Résultats et discussion", [
        "Modèle hybride FCFA + blockchain validé techniquement",
        "Forces : contextualisation Mali, inclusion Mobile Money, traçabilité optionnelle",
        "Limites : sandbox Mobile Money, testnet blockchain, pas d'audit externe",
        "Pas d'étude utilisateur terrain à grande échelle",
        "H1 et H3 confirmées ; H2 partiellement confirmée"
    ])

    # 13. Perspectives
    add_content_slide(prs, "Perspectives d'évolution", [
        "Court terme : API Mobile Money production, APK mobile finalisé, tests E2E",
        "Moyen terme : extension UEMOA (Sénégal, Côte d'Ivoire, Burkina), oracles Chainlink, pont FCFA/stablecoin",
        "Long terme : passage mainnet après audit, DAO, IA pour détection de fraude, partenariats télécoms"
    ])

    # 14. Démonstration
    add_demo_slide(prs, "Démonstration", [
        "Connexion client / prestataire",
        "Création mission + paiement Orange Money sandbox",
        "Candidature, acceptation, caution",
        "GPS, preuves, validation, KYC / litiges"
    ])

    # 15. Conclusion
    add_content_slide(prs, "Conclusion", [
        "BlockTask prouve la faisabilité d'une plateforme hybride adaptée au Mali",
        "Inclusion financière via Mobile Money + confiance institutionnelle via KYC NINA",
        "Traçabilité renforcée via smart contracts lorsque l'utilisateur le souhaite",
        "Ouverture vers l'UEMOA et la production"
    ])

    # 16. Remerciements
    add_title_slide(prs, "Merci de votre attention", "Questions ?", logo_path=LOGO_PATH)

    output_path = r"C:\Users\PC\Documents\blocktask_myself\rapport memoire\BlockTask_Soutenance.pptx"
    prs.save(output_path)
    print(f"Présentation générée : {output_path}")


if __name__ == "__main__":
    main()
